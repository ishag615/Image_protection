"""
Enhanced OCR and Text Extraction
Handles text extraction from multiple document formats using Tesseract and document libraries
"""

try:
    import pytesseract
except ImportError as e:
    raise ImportError(
        "pytesseract is required for OCR. "
        "Install with `pip install pytesseract` and ensure Tesseract OCR engine is installed on your system. "
        f"Original error: {e}"
    )

from PIL import Image
import cv2
import numpy as np
from pdf2image import convert_from_path
from docx import Document as DocxDocument
from pathlib import Path
from typing import List, Dict, Tuple, Any
import logging
import tempfile

logger = logging.getLogger(__name__)


class OCRProcessor:
    """
    Multi-format OCR and text extraction
    Handles images, PDFs, Word documents, and more
    """
    
    def __init__(self):
        """Initialize OCR processor"""
        self.supported_formats = {
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.tiff'],
            'pdf': ['.pdf'],
            'word': ['.docx', '.doc'],
            'powerpoint': ['.pptx', '.ppt']
        }
        logger.info("✓ OCR Processor initialized")
    
    def extract_from_file(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from any supported file format
        
        Args:
            file_path: Path to file
            
        Returns:
            Dictionary with extracted text and metadata
        """
        try:
            file_ext = Path(file_path).suffix.lower()
            
            # Route to appropriate extractor
            if file_ext in self.supported_formats['image']:
                return self._extract_from_image(file_path)
            elif file_ext == '.pdf':
                return self._extract_from_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self._extract_from_word(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self._extract_from_powerpoint(file_path)
            else:
                return {
                    'success': False,
                    'error': f'Unsupported file format: {file_ext}',
                    'text': '',
                    'metadata': {}
                }
        
        except Exception as e:
            logger.error(f"Error extracting from {file_path}: {e}")
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'metadata': {}
            }
    
    def _extract_from_image(self, image_path: str) -> Dict[str, Any]:
        """
        Extract text from image using Tesseract OCR
        Includes image preprocessing for better accuracy
        
        Args:
            image_path: Path to image
            
        Returns:
            Extraction result
        """
        try:
            img = Image.open(image_path)
            img_array = np.array(img)
            
            # Convert BGR to grayscale if needed
            if len(img_array.shape) == 3:
                if img_array.shape[2] == 4:  # RGBA
                    gray = cv2.cvtColor(img_array, cv2.COLOR_RGBA2GRAY)
                else:  # RGB
                    gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array
            
            # Apply preprocessing for better OCR
            processed = self._preprocess_image(gray)
            
            # Extract text with Tesseract
            full_text = pytesseract.image_to_string(processed)
            
            # Get detailed data (character bounding boxes, confidence)
            data = pytesseract.image_to_data(processed, output_type=pytesseract.Output.DICT)
            
            # Calculate confidence
            confidences = [int(c) for c in data['conf'] if c != '-1']
            avg_confidence = np.mean(confidences) if confidences else 0
            
            return {
                'success': True,
                'text': full_text,
                'confidence': round(avg_confidence, 2),
                'char_count': len(full_text),
                'word_count': len(full_text.split()),
                'metadata': {
                    'image_size': img.size,
                    'image_mode': img.mode,
                    'ocr_engine': 'Tesseract',
                    'preprocessing_applied': True
                }
            }
        
        except Exception as e:
            logger.error(f"Error extracting from image {image_path}: {e}")
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'metadata': {}
            }
    
    def _extract_from_pdf(self, pdf_path: str) -> Dict[str, Any]:
        """
        Extract text from PDF using pdf2image + OCR
        Also attempts to extract embedded text if available
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            Extraction result with per-page information
        """
        try:
            all_text = []
            all_confidence = []
            page_count = 0
            
            try:
                # Try to extract embedded text first (if PDF has selectable text)
                import PyPDF2
                with open(pdf_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text = page.extract_text()
                        if text.strip():
                            all_text.append(text)
                
                if all_text:
                    logger.info(f"Extracted embedded text from PDF: {len(all_text)} sections")
            except Exception as e:
                logger.debug(f"Could not extract embedded PDF text: {e}")
            
            # Convert PDF to images and OCR
            with tempfile.TemporaryDirectory() as temp_dir:
                images = convert_from_path(pdf_path, dpi=300)
                
                for page_num, image in enumerate(images, 1):
                    # Save temporarily
                    temp_path = Path(temp_dir) / f'page_{page_num}.png'
                    image.save(str(temp_path))
                    
                    # Extract text from this page
                    page_extraction = self._extract_from_image(str(temp_path))
                    
                    if page_extraction['success']:
                        all_text.append(page_extraction['text'])
                        if page_extraction.get('confidence'):
                            all_confidence.append(page_extraction['confidence'])
                    
                    page_count += 1
            
            combined_text = '\n--- PAGE BREAK ---\n'.join(all_text)
            avg_confidence = np.mean(all_confidence) if all_confidence else 0
            
            return {
                'success': True,
                'text': combined_text,
                'confidence': round(avg_confidence, 2),
                'page_count': page_count,
                'char_count': len(combined_text),
                'word_count': len(combined_text.split()),
                'metadata': {
                    'pdf_pages': page_count,
                    'extraction_method': 'OCR + Embedded',
                    'dpi': 300,
                    'ocr_engine': 'Tesseract'
                }
            }
        
        except Exception as e:
            logger.error(f"Error extracting from PDF {pdf_path}: {e}")
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'metadata': {}
            }
    
    def _extract_from_word(self, doc_path: str) -> Dict[str, Any]:
        """
        Extract text from Word document
        
        Args:
            doc_path: Path to .docx or .doc file
            
        Returns:
            Extraction result
        """
        try:
            # Handle .doc vs .docx
            if doc_path.endswith('.doc'):
                # For .doc files, might need libreoffice/pandoc conversion
                logger.warning(f"Legacy .doc file detected: {doc_path}. Using .docx extractor.")
                # For now, return error - real implementation would convert
                return {
                    'success': False,
                    'error': 'Legacy .doc format requires conversion. Use .docx instead.',
                    'text': '',
                    'metadata': {}
                }
            
            doc = DocxDocument(doc_path)
            
            # Extract paragraphs
            paragraphs = [para.text for para in doc.paragraphs]
            
            # Extract table content
            table_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_cells = [cell.text for cell in row.cells]
                    table_text.append(' | '.join(row_cells))
            
            # Combine all text
            all_text = paragraphs + table_text
            combined_text = '\n'.join(all_text)
            
            return {
                'success': True,
                'text': combined_text,
                'confidence': 1.0,  # Full confidence for extracted text (not OCR)
                'paragraph_count': len(paragraphs),
                'table_count': len(doc.tables),
                'char_count': len(combined_text),
                'word_count': len(combined_text.split()),
                'metadata': {
                    'document_type': 'Word',
                    'paragraphs': len(paragraphs),
                    'tables': len(doc.tables),
                    'extraction_method': 'Direct text extraction'
                }
            }
        
        except Exception as e:
            logger.error(f"Error extracting from Word document {doc_path}: {e}")
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'metadata': {}
            }
    
    def _extract_from_powerpoint(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract text from PowerPoint presentation
        
        Args:
            pptx_path: Path to .pptx file
            
        Returns:
            Extraction result
        """
        try:
            try:
                from pptx import Presentation
            except ImportError:
                return {
                    'success': False,
                    'error': 'python-pptx library required. Install it to process PowerPoint files.',
                    'text': '',
                    'metadata': {}
                }
            
            prs = Presentation(pptx_path)
            
            all_text = []
            slide_texts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                
                # Extract from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_content.append(shape.text)
                
                if slide_content:
                    slide_text = '\n'.join(slide_content)
                    all_text.append(f"[SLIDE {slide_num}]\n{slide_text}")
                    slide_texts.append(slide_text)
            
            combined_text = '\n--- SLIDE BREAK ---\n'.join(all_text)
            
            return {
                'success': True,
                'text': combined_text,
                'confidence': 1.0,  # Direct extraction
                'slide_count': len(prs.slides),
                'slides_with_text': len(slide_texts),
                'char_count': len(combined_text),
                'word_count': len(combined_text.split()),
                'metadata': {
                    'document_type': 'PowerPoint',
                    'slides': len(prs.slides),
                    'extraction_method': 'Direct text extraction'
                }
            }
        
        except Exception as e:
            logger.error(f"Error extracting from PowerPoint {pptx_path}: {e}")
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'metadata': {}
            }
    
    def _preprocess_image(self, gray_image: np.ndarray) -> np.ndarray:
        """
        Preprocess image for better OCR accuracy
        
        Args:
            gray_image: Grayscale image as numpy array
            
        Returns:
            Preprocessed image
        """
        # Apply CLAHE (Contrast Limited Adaptive Histogram Equalization)
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        enhanced = clahe.apply(gray_image)
        
        # Denoise - bilateral filter preserves edges
        denoised = cv2.bilateralFilter(enhanced, 9, 75, 75)
        
        # Sharpen text
        kernel = np.array([[-1, -1, -1],
                          [-1,  9, -1],
                          [-1, -1, -1]])
        sharpened = cv2.filter2D(denoised, -1, kernel)
        
        # Threshold to binary (helps with OCR)
        _, binary = cv2.threshold(sharpened, 150, 255, cv2.THRESH_BINARY)
        
        # Upscale for better recognition (tesseract works better with larger images)
        upscaled = cv2.resize(binary, None, fx=2, fy=2, interpolation=cv2.INTER_CUBIC)
        
        return upscaled
    
    def extract_text_regions(self, image_path: str) -> List[Dict[str, Any]]:
        """
        Extract text regions from image with bounding boxes
        Useful for targeted redaction
        
        Args:
            image_path: Path to image
            
        Returns:
            List of text regions with coordinates
        """
        try:
            img = Image.open(image_path)
            img_array = np.array(img)
            
            if len(img_array.shape) == 3:
                gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array
            
            processed = self._preprocess_image(gray)
            
            # Get detailed OCR data with bounding boxes
            data = pytesseract.image_to_data(processed, output_type=pytesseract.Output.DICT)
            
            regions = []
            for i in range(len(data['text'])):
                if int(data['conf'][i]) > 30:  # Confidence > 30%
                    regions.append({
                        'text': data['text'][i],
                        'confidence': int(data['conf'][i]),
                        'x': data['left'][i],
                        'y': data['top'][i],
                        'width': data['width'][i],
                        'height': data['height'][i],
                        'block_num': data['block_num'][i],
                        'line_num': data['line_num'][i],
                        'word_num': data['word_num'][i]
                    })
            
            return regions
        
        except Exception as e:
            logger.error(f"Error extracting text regions from {image_path}: {e}")
            return []
