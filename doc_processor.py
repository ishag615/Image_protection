from pathlib import Path
from PIL import Image, ImageDraw, ImageFilter
import io
import re

class DocumentProcessor:
    """Handle multiple document formats and image processing"""
    
    SUPPORTED_FORMATS = {
        'image': ['.jpg', '.jpeg', '.png', '.gif', '.webp'],
        'pdf': ['.pdf'],
        'word': ['.docx'],
        'powerpoint': ['.pptx']
    }
    
    def get_file_type(self, filepath):
        """Determine file type from extension"""
        ext = Path(filepath).suffix.lower()
        
        for file_type, extensions in self.SUPPORTED_FORMATS.items():
            if ext in extensions:
                return file_type
        
        raise ValueError(f"Unsupported file type: {ext}")
    
    def convert_to_images(self, filepath, file_type):
        """Convert document to images for analysis"""
        
        if file_type == 'image':
            # Already an image, just return the path
            return [filepath]
        
        elif file_type == 'pdf':
            # Convert PDF pages to images
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(filepath, dpi=200)
                image_paths = []
                
                for i, image in enumerate(images):
                    path = f'uploads/temp_pdf_page_{i}.png'
                    image.save(path)
                    image_paths.append(path)
                
                return image_paths
            except Exception as e:
                print(f"PDF conversion error: {e}")
                # Fallback: return original path
                return [filepath]
        
        elif file_type == 'word':
            # Extract images from DOCX
            try:
                from docx import Document
                doc = Document(filepath)
                image_paths = []
                
                # Extract inline images
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        try:
                            image_blob = rel.target_part.blob
                            path = f'uploads/temp_docx_image_{len(image_paths)}.png'
                            with open(path, 'wb') as f:
                                f.write(image_blob)
                            image_paths.append(path)
                        except:
                            pass
                
                # If no images found, create text rendering
                if not image_paths:
                    img = Image.new('RGB', (800, 1000), color='white')
                    draw = ImageDraw.Draw(img)
                    
                    text = ""
                    for para in doc.paragraphs:
                        text += para.text + "\n"
                    
                    # Draw text (basic rendering)
                    draw.text((20, 20), text[:500], fill='black')
                    
                    path = 'uploads/temp_docx_text.png'
                    img.save(path)
                    image_paths.append(path)
                
                return image_paths if image_paths else [filepath]
            except Exception as e:
                print(f"Word conversion error: {e}")
                return [filepath]
        
        elif file_type == 'powerpoint':
            # Extract images from PPTX
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                image_paths = []
                
                for slide_idx, slide in enumerate(prs.slides):
                    for shape_idx, shape in enumerate(slide.shapes):
                        if hasattr(shape, "image"):
                            try:
                                image_blob = shape.image.blob
                                path = f'uploads/temp_pptx_slide_{slide_idx}_{shape_idx}.png'
                                with open(path, 'wb') as f:
                                    f.write(image_blob)
                                image_paths.append(path)
                            except:
                                pass
                
                return image_paths if image_paths else [filepath]
            except Exception as e:
                print(f"PowerPoint conversion error: {e}")
                return [filepath]
        
        return [filepath]
    
    def redact_image(self, image_path):
        """Redact sensitive areas from image (basic version)"""
        try:
            img = Image.open(image_path).convert('RGB')
            
            # For now, just return the image as-is
            # Full redaction with Rekognition will be added later
            return img
        except Exception as e:
            print(f"Image redaction error: {e}")
            return None
    
    def extract_coordinates(self, instruction_text):
        """Extract x1,y1,x2,y2 from instruction"""
        numbers = re.findall(r'\d+', instruction_text)
        if len(numbers) >= 4:
            return tuple(map(int, numbers[:4]))
        return None